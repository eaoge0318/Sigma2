"""
PPTX Ingestor - python-pptx 版
每張投影片的文字抽出，輸出與 pdf_ingestor / docx_ingestor 相同格式
不支援 .ppt（舊格式），請另存為 .pptx 或 .pdf
"""
import uuid
import logging
from typing import List, Dict, Any

logger = logging.getLogger(__name__)

CHUNK_SIZE = 600
OVERLAP = 80


def ingest_pptx(raw: bytes, filename: str) -> Dict[str, Any]:
    """
    解析 PowerPoint (.pptx)，每張投影片為一個段落，切 chunk 後回傳統一格式
    """
    from pptx import Presentation
    from pptx.util import Pt
    import io

    doc_id = uuid.uuid4().hex[:12]
    prs = Presentation(io.BytesIO(raw))
    total_slides = len(prs.slides)

    text_chunks = []
    chunk_index = 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        lines = []

        # 標題優先
        if slide.shapes.title and slide.shapes.title.text.strip():
            lines.append(slide.shapes.title.text.strip())

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape == slide.shapes.title:
                continue
            for para in shape.text_frame.paragraphs:
                txt = para.text.strip()
                if txt:
                    lines.append(txt)

        slide_text = "\n".join(lines)
        if not slide_text.strip():
            continue

        for piece in _split_text(slide_text, CHUNK_SIZE, OVERLAP):
            text_chunks.append({
                "id": f"{doc_id}_s{slide_num}_c{chunk_index}",
                "doc_type": "text_chunk",
                "content": piece,
                "metadata": {
                    "source_name": filename,
                    "page": slide_num,
                    "chunk_index": chunk_index,
                    "file_type": "pptx",
                    "doc_id": doc_id,
                },
            })
            chunk_index += 1

    logger.info(f"[PPTX] '{filename}' 完成：slides={total_slides}，text_chunks={len(text_chunks)}")

    return {
        "doc_id": doc_id,
        "filename": filename,
        "text_chunks": text_chunks,
        "image_chunks": [],
        "manifest": {
            "total_pages": total_slides,
            "total_text_chunks": len(text_chunks),
            "total_images": 0,
        },
    }


def _split_text(text: str, chunk_size: int = 600, overlap: int = 80) -> List[str]:
    text = text.strip()
    if not text:
        return []
    if len(text) <= chunk_size:
        return [text]

    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        raw_chunk = text[start:end]
        for sep in ["。", "！", "？", "\n", ".", "!", "?"]:
            last = raw_chunk.rfind(sep)
            if last > chunk_size // 2:
                raw_chunk = raw_chunk[:last + 1]
                break
        advance = max(len(raw_chunk) - overlap, 1)
        chunk = raw_chunk.strip()
        if chunk:
            chunks.append(chunk)
        start += advance
    return chunks
